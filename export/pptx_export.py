# Copyright (c) 2026 Carlo Sicurini. All Rights Reserved.
# Metan.iQ - Biometano GHG Optimizer (DM 2022 / RED III)
# Proprietary and confidential. See LICENSE for terms.
# Commercial licensing: carlo.sicurini@gmail.com
﻿import io
import datetime
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False

def _fmt_num(value, decimals=0, suffix="", prefix=""):
    if value is None: return "N/D"
    try:
        if decimals == 0:
            formatted = f"{value:,.0f}".replace(",", "'")
        else:
            fmt_str = f"{{:,.{decimals}f}}"
            formatted = fmt_str.format(value).replace(",", "'").replace(".", ",")
            parts = formatted.split(",")
            if len(parts) > 1:
                formatted = parts[0].replace("'", ".") + "," + parts[1]
            else:
                formatted = parts[0].replace("'", ".")
        return prefix + formatted + suffix
    except:
        return prefix + str(value) + suffix

# Colori Stile Claude / Modern
C_BG = RGBColor(248, 249, 250)         # F8F9FA - Sfondo chiarissimo
C_TEXT_MAIN = RGBColor(30, 41, 59)     # 1E293B - Slate 800 (Testi principali)
C_TEXT_MUTED = RGBColor(100, 116, 139) # 64748B - Slate 500 (Testi secondari)
C_ACCENT = RGBColor(217, 119, 6)       # D97706 - Amber 600 (Accento principale)
C_SUCCESS = RGBColor(16, 185, 129)     # 10B981 - Emerald 500 (OK)
C_ERROR = RGBColor(244, 63, 94)        # F43F5E - Rose 500 (KO)
C_CARD_BG = RGBColor(255, 255, 255)    # FFFFFF - Sfondo Card

def build_metaniq_pptx(ctx: dict) -> io.BytesIO:
    if not _HAS_PPTX:
        raise ImportError("python-pptx non e' installato.")
        
    prs = Presentation()
    
    # Imposta proporzione 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    layout_blank = prs.slide_layouts[6]
    
    def add_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = C_BG

    def add_title(slide, text, x=0.5, y=0.4, size=32):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(12), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_MAIN
        # line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+0.6), Inches(1), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = C_ACCENT
        line.line.fill.background()

    def draw_kpi_card(slide, x, y, w, h, title, value, color_value=C_ACCENT, subtitle=""):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = RGBColor(226, 232, 240)
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.color.rgb = C_TEXT_MUTED
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(36)
        p2.font.color.rgb = color_value
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            p3 = tf.add_paragraph()
            p3.text = subtitle
            p3.font.size = Pt(11)
            p3.font.color.rgb = C_TEXT_MUTED
            p3.alignment = PP_ALIGN.CENTER

    # Estrazione Dati Base
    ghg_saving = ctx.get("saving_avg") or 0.0
    mwh = ctx.get("tot_mwh_basis") or 0.0
    rev = ctx.get("tot_revenue") or 0.0
    taglia = ctx.get("plant_net_smch") or 300.0
    aux = ctx.get("aux_factor") or 1.29
    valid_months = ctx.get("valid_months") or 0
    compliant = ghg_saving >= 80 and valid_months == 12
    
    # Dati BP
    bp = ctx.get("bp_result") or {}
    opex_list = bp.get("opex", []) if isinstance(bp, dict) else []
    avg_opex = sum(opex_list)/len(opex_list) if opex_list else 0.0
    capex = bp.get("capex_totale", 0.0) if isinstance(bp, dict) else 0.0
    ebitda_list = bp.get("ebitda", []) if isinstance(bp, dict) else []
    avg_ebitda = sum(ebitda_list)/len(ebitda_list) if ebitda_list else 0.0
    payback = bp.get("payback_anno", "N/D") if isinstance(bp, dict) else "N/D"

    # ==========================================
    # SLIDE 1: Copertina
    # ==========================================
    s1 = prs.slides.add_slide(layout_blank)
    add_background(s1)
    
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.5), Inches(0.15), Inches(2.0))
    line.fill.solid()
    line.fill.fore_color.rgb = C_ACCENT
    line.line.fill.background()
    
    tx = s1.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(10), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Metan.iQ Simulator"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_MAIN
    
    p2 = tf.add_paragraph()
    p2.text = "Business Plan & Analisi Sostenibilita' RED III"
    p2.font.size = Pt(24)
    p2.font.color.rgb = C_TEXT_MUTED
    
    p3 = tf.add_paragraph()
    p3.text = f"\nData di estrazione: {datetime.datetime.now().strftime('%d/%m/%Y')}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_TEXT_MUTED
    
    # ==========================================
    # SLIDE 2: Executive Summary
    # ==========================================
    s2 = prs.slides.add_slide(layout_blank)
    add_background(s2)
    add_title(s2, "Executive Summary")
    
    draw_kpi_card(s2, 0.5, 1.8, 3.8, 2.2, "SAVING GHG MEDIO", _fmt_num(ghg_saving, 1, '%'), color_value=C_SUCCESS if compliant else C_ERROR, subtitle="Soglia Minima: 80%")
    draw_kpi_card(s2, 4.75, 1.8, 3.8, 2.2, "CONFORMITA' RED III", "IDONEO" if compliant else "NON IDONEO", color_value=C_SUCCESS if compliant else C_ERROR, subtitle=f"Mesi conformi: {valid_months}/12")
    draw_kpi_card(s2, 9.0, 1.8, 3.8, 2.2, "ENERGIA IMMESSA", _fmt_num(mwh, 0, ' MWh/anno'), subtitle=f"Da {_fmt_num(taglia,0)} Sm3/h netti autorizzati")
    
    draw_kpi_card(s2, 0.5, 4.3, 5.9, 2.0, "RICAVI MEDI ANNUI", _fmt_num(rev, 0, ' €', ' '), subtitle="Include tariffa GSE e mercato")
    draw_kpi_card(s2, 6.9, 4.3, 5.9, 2.0, "EBITDA MEDIO", _fmt_num(avg_ebitda, 0, ' €', ' '), subtitle="Margine operativo lordo post-OPEX")
    # ==========================================
    # SLIDE 3: Mix Biomasse e Sostenibilità
    # ==========================================
    s3 = prs.slides.add_slide(layout_blank)
    add_background(s3)
    add_title(s3, "Configurazione Impianto & Mix Biomasse")
    
    feedstock_rows = ctx.get("revenue_rows") or []
    rows = len(feedstock_rows) + 1
    cols = 3
    if rows > 1:
        table_shape = s3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(12.3), Inches(rows * 0.6))
        table = table_shape.table
        
        headers = ["Biomassa", "Fabbisogno Annuo", "MWh Producibili"]
        for c in range(cols):
            cell = table.cell(0, c)
            cell.text = headers[c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_TEXT_MAIN
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            p.font.size = Pt(16)
            
        for r, (name, data) in enumerate(feedstock_rows):
            table.cell(r+1, 0).text = name
            table.cell(r+1, 1).text = _fmt_num(data.get("t", 0), 0, " tonnellate")
            table.cell(r+1, 2).text = _fmt_num(data.get("mwh", 0), 0, " MWh")
            for c in range(cols):
                p = table.cell(r+1, c).text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.color.rgb = C_TEXT_MAIN
                
    tx = s3.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Il fabbisogno e' dimensionato considerando gli ausiliari ({_fmt_num((aux-1)*100, 0, '%')}). Base calcolo GHG: {_fmt_num(taglia*aux,0)} Sm3/h lordi equivalenti."
    p.font.size = Pt(16)
    p.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 4: Conto Economico & Ricavi
    # ==========================================
    s4 = prs.slides.add_slide(layout_blank)
    add_background(s4)
    add_title(s4, "Struttura dei Ricavi & Base Tariffaria")
    
    tariffa = ctx.get("tariffa_media_ponderata") or 0.0
    
    draw_kpi_card(s4, 0.5, 1.8, 5.9, 2.5, "TARIFFA MEDIA PONDERATA", _fmt_num(tariffa, 2, ' €/MWh'), subtitle="Risultato base d'asta post-ribasso")
    draw_kpi_card(s4, 6.9, 1.8, 5.9, 2.5, "RICAVO TOTALE ANNUO", _fmt_num(rev, 0, ' €'), subtitle="Proiezione media")

    tx = s4.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(2))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.add_paragraph().text = f"Il valore cumulato dei ricavi sui 15 anni di vita utile del progetto si attesta intorno a {_fmt_num(rev * 15, 0, ' €')}."
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = C_TEXT_MAIN
    tf.paragraphs[0].font.bold = True

    # ==========================================
    # SLIDE 7: Piano di Produzione
    # ==========================================
    s7 = prs.slides.add_slide(layout_blank)
    add_background(s7)
    add_title(s7, "Piano di Produzione e Regolazione Mensile")
    
    tx = s7.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(4))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.add_paragraph().text = "Il Metan.iQ Solver verifica mensilmente il mix in ingresso."
    tf.add_paragraph().text = f"Mesi Conformi dal punto di vista ambientale: {valid_months} su 12."
    tf.add_paragraph().text = ""
    tf.add_paragraph().text = "Un mese non conforme NON preclude l'accesso agli incentivi qualora la media annuale compensi il deficit emissivo."
    
    for i, para in enumerate(tf.paragraphs):
        para.font.size = Pt(20)
        para.font.color.rgb = C_TEXT_MAIN

    # ==========================================
    # SLIDE 8: Conclusioni
    # ==========================================
    s8 = prs.slides.add_slide(layout_blank)
    add_background(s8)
    add_title(s8, "Conclusioni e Valutazione")
    
    tx = s8.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(4))
    tf = tx.text_frame
    tf.word_wrap = True
    
    if compliant:
        conc_text = "FATTIBILITÀ CONFERMATA\n\nIl mix di biomasse analizzato rispetta pienamente la Direttiva RED III."
        tf.paragraphs[0].font.color.rgb = C_SUCCESS
    else:
        conc_text = "ALLERTA RED III\n\nL'impianto supera i limiti di emissione.\nÈ obbligatorio ricalibrare il mix."
        tf.paragraphs[0].font.color.rgb = C_ERROR

    for i, line in enumerate(conc_text.split('\n')):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(24)
        if i > 0:
            p.font.color.rgb = C_TEXT_MAIN
            
    # Footer in tutte le slide
    for slide in prs.slides:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.5))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Metan.iQ Simulator | Riservato"
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED
        p.alignment = PP_ALIGN.RIGHT

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
